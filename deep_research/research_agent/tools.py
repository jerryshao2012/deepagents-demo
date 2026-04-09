"""Research Tools.

This module provides search and content processing utilities for the research agent,
using Tavily for URL discovery and fetching full webpage content.
"""

from __future__ import annotations

import datetime
import json
import os
import random
import re
from json import dumps as json_dumps
from pathlib import Path

import httpx
import jsonschema
import pymupdf4llm
import pypdf
import requests
from deepagents.backends.utils import create_file_data
from docx import Document
from dotenv import load_dotenv
from langchain_core.tools import InjectedToolArg, tool
from langgraph.prebuilt import InjectedState
from markdownify import markdownify
from openpyxl import load_workbook
from pptx import Presentation
from tavily import TavilyClient
from typing_extensions import Annotated, Literal

from research_agent.targets import get_target_definition
from utils import get_ssl_verify_config

load_dotenv()

verify_ssl = get_ssl_verify_config()
tavily_session = requests.Session()
tavily_session.verify = verify_ssl
tavily_client = TavilyClient(session=tavily_session)
# Global in‑memory cache for folder listings (path → list of Path objects)
_folder_listing_cache: dict[str, list[Path]] = {}

# Limit recursion depth for recursive glob to avoid walking deep trees
MAX_GLOB_DEPTH = 3
SUPPORTED_DOC_SUFFIXES = {".pdf", ".txt", ".md", ".docx", ".pptx", ".xlsx"}

REPORTS_OUTPUT_FOLDER = "./output"
MAX_FILES_TO_READ = 20
MAX_TOTAL_SIZE_MB = 50


def _normalize_path_for_filesystem_tools(path_str: str) -> str:
    """Normalize paths for cross-platform compatibility with deepagents filesystem tools.
    
    Deepagents filesystem tools (glob, ls, etc.) expect paths relative to the working directory.
    This function ensures paths start with './' instead of '/' for proper resolution on all platforms.
    
    Args:
        path_str: The path string to normalize
        
    Returns:
        Normalized path string with proper relative prefix
    """
    if not path_str:
        return path_str

    # Convert Windows backslashes to forward slashes for consistency
    normalized = path_str.replace('\\', '/')

    # If path starts with '/', it's being treated as absolute from root
    # Convert to relative path by adding './' prefix
    if normalized.startswith('/') and not normalized.startswith('./'):
        normalized = './' + normalized.lstrip('/')

    return normalized


def _run_tavily_search(query: str, max_results: int, topic: str, timeout: float = 60.0) -> dict:
    api_key = os.getenv("TAVILY_API_KEY")
    if not api_key:
        raise ValueError("TAVILY_API_KEY is not set")

    payload = {
        "api_key": api_key,
        "query": query,
        "max_results": max_results,
        "topic": topic,
        "include_answer": False,
        "include_raw_content": False,
    }
    response = tavily_client.session.post(
        f"{tavily_client.base_url}/search",
        data=json_dumps(payload),
        timeout=min(timeout, 120),
        verify=verify_ssl,
        headers={"Content-Type": "application/json"},
    )
    response.raise_for_status()
    response_dict = response.json()
    if not isinstance(response_dict, dict):
        return {"results": []}
    response_dict.setdefault("results", [])
    return response_dict


def fetch_webpage_content(url: str, timeout: float = 10.0) -> str:
    """Fetch and convert webpage content to markdown."""
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        )
    }

    try:
        response = httpx.get(
            url=url,
            headers=headers,
            timeout=timeout,
            verify=verify_ssl,
        )
        response.raise_for_status()
        return markdownify(response.text)
    except Exception as exc:
        return f"Error fetching content from {url}: {exc}"


def _extract_pdf_text(file_path: Path) -> str:
    """Extract PDF content as markdown without ML model downloads.\n\n    Returns:\n        str: Extracted content from the PDF file.\n    """
    try:
        print("Use PyMuPDF4LLM for PDF markdown extraction.")

        markdown_content = pymupdf4llm.to_markdown(str(file_path))
        if isinstance(markdown_content, list):
            # Convert each dictionary to a string representation
            return "\n\n".join(str(item) for item in markdown_content)
        if markdown_content.strip():
            return markdown_content
    except Exception as e:
        print(f"PyMuPDF4LLM PDF extraction failed: {e}")
        # Fallback to pypdf if markdown extraction fails
        try:
            print("Falling back to pypdf for PDF text extraction.")
            reader = pypdf.PdfReader(file_path)
            page_texts: list[str] = []
            for index, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    page_texts.append(f"## Page {index}:\n\n{text}")
            return "\n\n".join(page_texts)
        except Exception as e:
            return f"Error extracting PDF text: {e}"
    return ""


def _extract_text_file(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8")


def _extract_docx_text(file_path: Path) -> str:
    document = Document(str(file_path))
    paragraphs = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text and paragraph.text.strip()
    ]

    # Many documents are table-based; include table cells so content is not silently missed.
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                call_paragraph = [
                    paragraph.text.strip()
                    for paragraph in cell.paragraphs
                    if paragraph.text and paragraph.text.strip()
                ]
                if call_paragraph:
                    paragraphs.extend(call_paragraph)

    return "\n".join(paragraphs)


def _extract_pptx_text(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    slide_sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide:
            notes = []
            for shape in slide.notes_slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    notes.append(text.strip())
            notes_text = "\n".join(notes)
        if notes_text:
            parts.append(f"Speaker Notes:\n{notes_text}")
        slide_sections.append("\n".join(parts))

    return "\n\n".join(slide_sections)


def _extract_xlsx_text(file_path: Path) -> str:
    workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    sections: list[str] = []
    try:
        for worksheet in workbook.worksheets:
            rows: list[str] = []
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if values:
                    rows.append(" | ".join(values))
            body = "\n".join(rows) if rows else "(empty sheet)"
            sections.append(f"Sheet: {worksheet.title}\n{body}")
    finally:
        workbook.close()

    return "\n\n".join(sections)


def _save_extracted_content(original_file_path: Path, content: str, output_folder: Path | None = None) -> str:
    """Save extracted content to the output folder with appropriate extension.

    Args:
        original_file_path: Path to the original document
        content: Extracted text/markdown content
        output_folder: Optional subfolder within REPORTS_OUTPUT_FOLDER to save to

    Returns:
        str: Path to the saved file
    """
    if output_folder:
        output_dir = output_folder
    else:
        output_dir = Path(REPORTS_OUTPUT_FOLDER)

    output_dir.mkdir(parents=True, exist_ok=True)

    # Use original filename (without extension) + new extension
    file_path = _get_extracted_path(original_file_path, output_dir)
    file_path.parent.mkdir(parents=True, exist_ok=True)

    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)

    return str(file_path)


def _extract_supported_document(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf_text(file_path)
    if suffix in {".txt", ".md"}:
        return _extract_text_file(file_path)
    if suffix == ".docx":
        return _extract_docx_text(file_path)
    if suffix == ".pptx":
        return _extract_pptx_text(file_path)
    if suffix == ".xlsx":
        return _extract_xlsx_text(file_path)
    raise ValueError(f"Unsupported document type: {suffix}")


@tool(parse_docstring=True)
def tavily_search(
        query: str,
        max_results: Annotated[int, InjectedToolArg] = 1,
        topic: Annotated[Literal["general", "news", "finance"], InjectedToolArg] = "general",
        state: Annotated[dict, InjectedState] = None,
) -> str:
    """Search the web for information on a given query.

    Uses Tavily to discover relevant URLs, then fetches and returns full webpage content as markdown.

    Args:
        query: Search query to execute
        max_results: Maximum number of results to return (default: 1)
        topic: Topic filter - 'general', 'news', or 'finance' (default: 'general')
        state: LangGraph state

    Returns:
        Formatted search results with full webpage content
    """
    if state:
        messages = state.get("messages", [])
        for msg in messages:
            content = getattr(msg, "content", "") if not isinstance(msg, dict) else msg.get("content", "")
            if isinstance(content, str) and "Do NOT use web search" in content:
                return "Note: Web search is disabled for this research task. Please use local documents or internal knowledge only."

    try:
        search_results = _run_tavily_search(
            query=query,
            max_results=max_results,
            topic=topic,
        )
    except requests.exceptions.HTTPError as exc:
        status_code = exc.response.status_code if exc.response is not None else None
        if status_code == 401:
            return (
                "Tavily authentication failed (401 Unauthorized)."
                "Set a valid TAVILY_API_KEY environment variable and retry"
            )
        return f"Tavily request failed with HTTP {status_code}: {exc}"
    except ValueError as exc:
        return str(exc)
    except Exception as exc:
        return f"Tavily search failed: {exc}"

    result_texts = []
    for result in search_results.get("results", []):
        url = result["url"]
        title = result["title"]
        content = fetch_webpage_content(url)
        result_text = f"""## {title}
**URL:** {url}

{content}

---
"""
        result_texts.append(result_text)

    return f"""🔍 Found {len(result_texts)} result(s) for '{query}':

{chr(10).join(result_texts)}"""


@tool(parse_docstring=True)
def think_tool(reflection: str) -> str:
    """Tool for strategic reflection on research progress and decision-making.

    Use this tool after each search to analyze results and plan next steps systematically.
    This creates a deliberate pause in the research workflow for quality decision-making.

    Args:
        reflection: Detailed reflection on research progress, findings, gaps, and next steps

    Returns:
        Confirmation that reflection was recorded for decision-making
    """
    return f"Reflection recorded: {reflection}"


def _get_extracted_path(file_path: Path, output_folder: Path) -> Path:
    """Get the target path for an extracted file."""
    suffix = file_path.suffix.lower()
    if suffix in {".pdf", ".md", ".docx", ".pptx"}:
        new_extension = ".md"
    else:
        new_extension = ".txt"

    new_filename = f"{file_path.name}_extracted{new_extension}"
    return output_folder / "extracted" / new_filename


def _resolve_doc_output_subfolder(folder: Path) -> Path:
    configured_output = Path(os.environ.get("OUTPUT_FOLDER", REPORTS_OUTPUT_FOLDER))
    if configured_output.name == folder.name:
        return configured_output
    if configured_output == Path(REPORTS_OUTPUT_FOLDER):
        return configured_output / folder.name
    return configured_output


@tool(parse_docstring=True)
def read_doc_folder(
        folder_path: str,
        specific_files: list[str] | None = None,
        state: Annotated[dict, InjectedState] = None,
) -> str:
    """Read and extract text from supported documents in a given folder.

    Use this tool when you need to research from local documents instead of or in addition
    to web search. Supported file types are PDF, text, Markdown, Word, PowerPoint, and Excel.

    If the folder contains a large number of files or the total size is very large,
    this tool will return a summary of the contents instead of all text.
    You can then use the `specific_files` argument to read particular documents of interest.

    Args:
        folder_path: The absolute or relative path to the folder containing document files.
        specific_files: Optional list of filenames within the folder to read specifically.
            If provided, only these files will be processed, bypassing general limits.
        state: LangGraph state (injected automatically, do not supply).

    Returns:
        Extracted text from supported documents, a summary for large folders, or an error message.
    """
    configured_doc_folder: str | None = None
    if state and isinstance(state, dict):
        configured_doc_folder = state.get("doc_folder")

    # No doc_folder configured → block all filesystem access.
    # The tool is only meaningful when an explicit folder has been provided via --doc-folder.
    if not configured_doc_folder:
        return (
            "Error: No document folder has been configured for this research task. "
            "Pass --doc-folder <path> when invoking the CLI to enable local document reading. "
            "Do NOT attempt to read from any other filesystem path."
        )

    allowed_root = Path(configured_doc_folder).resolve()

    # Resolve the requested path and check it is within the allowed root.
    folder = Path(folder_path).resolve()
    try:
        folder.relative_to(allowed_root)
        # Path is within allowed_root — use it as-is.
    except ValueError:
        # Requested path is outside the configured doc_folder → redirect to the root.
        print(
            f"[read_doc_folder] Redirecting '{folder_path}' → '{allowed_root}' "
            f"(only the configured doc_folder is permitted)."
        )
        folder = allowed_root

    if not folder.exists():
        return f"Error: Folder '{folder}' does not exist."
    if not folder.is_dir():
        return f"Error: '{folder}' is not a directory."

    # Use a set for faster lookup if specific_files is provided
    specific_set = set(specific_files) if specific_files else None

    # ---------- Cached folder listing ----------
    cache_key = str(folder.resolve())
    if cache_key in _folder_listing_cache:
        supported_files = _folder_listing_cache[cache_key]
    else:
        all_candidates: list[Path] = []
        for file_path in folder.rglob("*"):
            # Enforce max recursion depth to avoid walking deep trees
            if len(file_path.relative_to(folder).parts) > MAX_GLOB_DEPTH:
                continue
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_DOC_SUFFIXES:
                all_candidates.append(file_path)
        supported_files = sorted(all_candidates)
        _folder_listing_cache[cache_key] = supported_files
    # ------------------------------------------

    if not supported_files:
        return (
            f"No supported document files found in {folder_path}. "
            "Supported types: .pdf, .txt, .md, .docx, .pptx, .xlsx."
        )

    # If specific files are requested, narrow the list down
    if specific_set:
        files_to_process = [f for f in supported_files if f.name in specific_set]
        if not files_to_process:
            return f"None of the requested files were found in {folder_path}. Available: {', '.join(f.name for f in supported_files[:10])}..."
    else:
        # Check limits for general folder reading
        total_files = len(supported_files)
        # Total size in MB (using lstat for speed)
        total_size_mb = sum(f.lstat().st_size for f in supported_files) / (1024 * 1024)

        if total_files > MAX_FILES_TO_READ or total_size_mb > MAX_TOTAL_SIZE_MB:
            # Randomly select diverse sample from the entire file collection
            # This ensures better coverage across all documents rather than just top ones
            # Bound sample size by both file count limit and estimated size limit
            avg_size_mb = total_size_mb / total_files if total_files > 0 else 0
            max_files_by_size = max(1, int(MAX_TOTAL_SIZE_MB / avg_size_mb)) if avg_size_mb > 0 else MAX_FILES_TO_READ
            sample_size = min(MAX_FILES_TO_READ, total_files, max_files_by_size)

            auto_sample = [f.name for f in random.sample(supported_files, sample_size)]

            # Show at most 60 files in the full listing so the context isn't blown out
            preview_list = "\n".join(
                f"- {f.name} ({f.lstat().st_size / 1024:.1f} KB)"
                for f in supported_files[:60]
            )
            if total_files > 60:
                preview_list += f"\n... and {total_files - 60} more files (not shown)."

            auto_sample_str = ", ".join(f'"{n}"' for n in auto_sample)

            return (
                f"TOOL RESULT — folder too large to read all at once: {total_files} files, {total_size_mb:.1f} MB "
                f"(limits: {MAX_FILES_TO_READ} files / {MAX_TOTAL_SIZE_MB} MB).\n\n"
                "ACTION REQUIRED — do NOT ask the user for confirmation. You MUST immediately:\n"
                f"1. Call read_doc_folder again on '{folder_path}' with specific_files set to the auto-sample below.\n"
                "2. Continue research using those documents.\n\n"
                f"Pre-built diverse auto-sample ({len(auto_sample)} files, evenly spread across the directory):\n"
                f"[{auto_sample_str}]\n\n"
                f"Full file listing (first 60 of {total_files}):\n{preview_list}"
            )
        files_to_process = supported_files

    extracted_text: list[str] = []
    processed_files: list[str] = []
    failed_files: list[str] = []
    output_subfolder = _resolve_doc_output_subfolder(folder)

    for file_path in files_to_process:
        target_path = _get_extracted_path(file_path, output_subfolder)
        if target_path.exists():
            print(f"Skipping {file_path.name}, already extracted to {target_path}")
            try:
                content = target_path.read_text(encoding="utf-8")
                processed_files.append(f"{file_path.name} (skipped, loaded from {target_path})")
                extracted_text.append(f"--- Content of {file_path.name} (from cache) ---\n{content}\n")
                continue
            except Exception as exc:
                print(f"Failed to read existing extract {target_path}: {exc}. Re-extracting...")

        print(f"Processing document: {file_path.name}...")
        try:
            content = _extract_supported_document(file_path)
            # Save the extracted content as well
            saved_path = _save_extracted_content(file_path, content, output_folder=output_subfolder)
            processed_files.append(f"{file_path.name} (saved to {saved_path})")
            extracted_text.append(f"--- Content of {file_path.name} ---\n{content}\n")
        except Exception as exc:
            failed_files.append(file_path.name)
            extracted_text.append(f"--- Error reading {file_path.name}: {exc} ---\n")

    summary_lines = [f"Processed {len(processed_files)}/{len(files_to_process)} supported file(s) from {folder}."]
    if processed_files:
        summary_lines.append(f"Files processed: {', '.join(processed_files)}")
    if failed_files:
        summary_lines.append(f"Files failed: {', '.join(failed_files)}")
    print("\n".join(summary_lines))
    return "\n".join(summary_lines + [""] + extracted_text)


def save_research_report(report_title: str, content: str) -> str:
    """Save a research report to the output folder.

    Args:
        report_title: Title of the report (used for filename)
        content: Content of the report

    Returns:
        str: Path to the saved report file
    """
    # For output subfolder
    output_subfolder = Path(os.environ.get("OUTPUT_FOLDER", REPORTS_OUTPUT_FOLDER))
    output_subfolder.mkdir(parents=True, exist_ok=True)

    # Clean up the report title to make it filename-friendly
    safe_title = re.sub(r"[^a-zA-Z0-9_\- ]", "", report_title).strip()[:100]
    safe_title = safe_title.replace(" ", "_")

    # Create filename with timestamp
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{timestamp}_{safe_title}.md"

    file_path = output_subfolder / filename

    with open(file_path, 'w', encoding='utf-8') as f:
        f.write(content)

    return str(file_path)


def _coerce_integers(value, schema):
    if isinstance(value, dict):
        props = schema.get('properties', {})
        return {k: _coerce_integers(v, props.get(k, {})) for k, v in value.items()}
    if isinstance(value, list):
        item_schema = schema.get('items', {})
        return [_coerce_integers(item, item_schema) for item in value]
    if schema.get("type") == "integer" and isinstance(value, float) and value.is_integer():
        return int(value)
    return value


def _resolve_path(path: str, context: dict[str, object]):
    if path == "item":
        return context.get("item")
    target = context["root"]
    if path.startswith("item."):
        target = context.get("item", {})
        path = path[5:]
    elif path.startswith("root."):
        path = path[5:]

    if not path:
        return target

    current = target
    for part in path.split("."):
        if isinstance(current, dict):
            current = current.get(part)
        else:
            return None
    return current


def _evaluate_expression(expression: str, context: dict[str, object]) -> str:
    expression = expression.strip()
    if expression == "index":
        return str(context.get("index", ""))
    if expression.startswith("sum(") and expression.endswith(")"):
        inner = expression[4:-1]
        array_path, _, field_name = inner.partition("[].")
        values = _resolve_path(array_path, context)
        if not isinstance(values, list) or not field_name:
            return ""
        total = 0
        for value in values:
            if isinstance(value, dict):
                number = value.get(field_name)
                if isinstance(number, (int, float)):
                    total += number
        return str(int(total) if isinstance(total, float) and total.is_integer() else total)

    value = _resolve_path(expression, context)
    if value is None:
        return ""
    return str(value)


def _interpolate_text(template: str, context: dict[str, object]) -> str:
    result = template
    for match in re.finditer(r"\{([^{}]+)}", template):
        expression = match.group(1)
        result = result.replace(match.group(0), _evaluate_expression(expression, context))
    return result


def _render_blocks(spec: list[dict[str, object]], context: dict[str, object]) -> list[str]:
    output: list[str] = []
    for block in spec:
        block_type = block.get("type")
        if block_type == "heading":
            level_value = block.get("level", 1)
            level = int(str(level_value))
            value = _interpolate_text(str(block.get("value", "")), context)
            output.append(f"{'#' * level} {value}".rstrip())
        elif block_type == "text":
            output.append(_interpolate_text(str(block.get("value", "")), context))
        elif block_type == "separator":
            output.append(str(block.get("value", "---")))
        elif block_type == "bullet_list":
            values = _resolve_path(str(block.get("path", "")), context)
            if isinstance(values, list):
                for value in values:
                    output.append(f"- {value}")
        elif block_type == "repeat":
            items = _resolve_path(str(block.get("path", "")), context)
            body = block.get("body", [])
            if isinstance(items, list) and isinstance(body, list):
                for index, item in enumerate(items, start=1):
                    child_context = {"root": context["root"], "item": item, "index": index}
                    output.extend(_render_blocks(body, child_context))
        elif block_type == "if_present":
            value = _resolve_path(str(block.get("path", "")), context)
            body = block.get("body", [])
            if value and isinstance(body, list):
                output.extend(_render_blocks(body, context))
        else:
            raise ValueError(f"Unsupported render block type: {block_type}")
    return output


def _fill_defaults(target_id: str, payload: dict) -> dict:
    """Supply sensible default values for optional metadata fields before schema validation.

    This prevents the LLM from failing validation on non-critical header fields and
    eliminates the need to ask the user for things like dataset_name or topic.
    """
    definition = get_target_definition(target_id)
    defaults = definition.get("defaults", [])
    if not isinstance(defaults, list):
        # Fallback if someone used a dict for some reason
        return payload

    # Pre-extract items if needed by rules
    # This helps with legacy support or unified access to 'questions', 'slides', etc.
    items = payload.get("items") or payload.get("questions") or payload.get("slides") or []

    for rule in defaults:
        field = rule.get("field")
        if not field:
            continue

        condition = rule.get("if_null", False)
        if condition and payload.get(field):
            continue

        expr = rule.get("value")
        if not expr:
            continue

        # Evaluate simple expressions or built-ins
        if expr == "items":
            payload[field] = items
        elif expr.startswith("first_of:"):
            fields = expr[9:].split(",")
            for f in fields:
                val = payload.get(f.strip())
                if val:
                    if isinstance(val, list) and len(val) > 0:
                        payload[field] = val[0]
                    else:
                        payload[field] = val
                    break
        elif expr == "collect_unique:coverage_area":
            seen: list[str] = []
            for item in items:
                area = item.get("coverage_area", "")
                if area and area not in seen:
                    seen.append(area)
            payload[field] = seen or ["General"]
        elif expr == "derive_dataset_name":
            domain = payload.get("domain", "")
            areas = payload.get("coverage_areas") or []
            if domain:
                payload[field] = f"{domain} Q&A Draft Set"
            elif areas:
                payload[field] = f"{areas[0]} Q&A Draft Set"
            else:
                payload[field] = "Golden Dataset Draft Set"
        elif expr == "derive_topic":
            first_val = ""
            if items:
                # Try common fields for topic derivation
                first_val = items[0].get("question") or items[0].get("title") or ""
            payload[field] = (first_val[:80] if first_val else "General")
        elif expr == "derive_objective":
            topic = payload.get("topic", "the subject")
            payload[field] = f"Assess knowledge and practical experience related to {topic}."
        elif expr == "dataset_size_calc":
            payload[field] = max(50, len(items) * 4)
        elif expr == "ensure_item_ids":
            for idx, item in enumerate(items, start=1):
                if not item.get("id"):
                    item["id"] = str(idx)
        elif expr == "ensure_item_content":
            for item in items:
                if "content" not in item:
                    item["content"] = ""
        else:
            # Simple literal string or fallback
            payload[field] = expr

    return payload


def _render_payload(template: str, payload, render_spec: list[dict[str, object]]) -> str:
    if template == "markdown_blocks":
        context = {"root": payload, "item": payload, "index": 1}
        rendered = [block for block in _render_blocks(render_spec, context) if block != ""]
        return "\n\n".join(rendered).strip() + "\n"
    raise ValueError(f"Unsupported render template: {template}")


def _normalize_legacy_target_payload(target_id: str, payload: dict) -> dict:
    """Apply narrow target-specific compatibility fixes before validation."""
    if target_id != "study-slides":
        return payload

    slides = payload.get("slides")
    if not isinstance(slides, list):
        return payload

    normalized_slides: list[dict] = []
    for slide in slides:
        if not isinstance(slide, dict):
            normalized_slides.append(slide)
            continue

        normalized_slide = dict(slide)

        # Older prompts sometimes used `content` instead of `bullets`.
        if "bullets" not in normalized_slide and "content" in normalized_slide:
            content = normalized_slide.get("content")
            if isinstance(content, list):
                normalized_slide["bullets"] = [str(item) for item in content]
            elif isinstance(content, str):
                normalized_slide["bullets"] = [content]

        # Some generations add metadata the schema does not allow.
        normalized_slide.pop("content", None)
        normalized_slide.pop("slide_number", None)

        normalized_slides.append(normalized_slide)

    payload = dict(payload)
    payload["slides"] = normalized_slides
    return payload


def _prepare_validated_payload(
        target_id: str, payload_json: str | dict
) -> tuple[dict | None, dict | None, str | None]:
    """Parse JSON, apply defaults, validate against schema. Returns (definition, payload, error_message)."""
    try:
        definition = get_target_definition(target_id)
    except ValueError as exc:
        return None, None, str(exc)

    if not definition.get("schema"):
        return None, None, f"ERROR: Target '{target_id}' is an unstructured target. Do NOT use `render_target_output`! You must formulate your response directly as markdown and write it to the final report file or output it."

    if isinstance(payload_json, dict):
        payload = payload_json
    else:
        try:
            payload = json.loads(payload_json)
        except json.JSONDecodeError as exc:
            return None, None, f"Invalid JSON payload: {exc}"

    payload = _normalize_legacy_target_payload(target_id, payload)

    payload = _fill_defaults(target_id, payload)
    payload = _coerce_integers(payload, definition["schema"])

    try:
        jsonschema.validate(instance=payload, schema=definition["schema"])
    except jsonschema.ValidationError as exc:
        return None, None, f"Schema validation failed for target '{target_id}': {exc.message}"

    return definition, payload, None


@tool(parse_docstring=True)
def render_target_output(
        target_id: str,
        payload_json: str | dict,
        state: Annotated[dict, InjectedState] = None,
) -> str:
    """Render structured target output using a reusable target definition.

    Use this tool ONLY for structured output targets (targets with a JSON schema).
    DO NOT use this tool for 'Unstructured Markdown Document' targets.
    Provide the target id and a JSON payload that matches the selected target schema exactly.
    The payload may be either a JSON object string or a dict-like JSON object.
    NEVER put raw markdown into payload_json.

    Args:
        target_id: The target definition id to use for validation and rendering.
        payload_json: A JSON object string or dict matching the target schema.
        state: LangGraph state

    Returns:
        Rendered markdown output or a validation error message.
    """
    definition, payload, err = _prepare_validated_payload(target_id, payload_json)
    if err:
        return err

    return _render_payload(definition["render"]["template"], payload, definition["render"]["spec"])


@tool(parse_docstring=True)
def finalize_golden_dataset_output(
        payload_json: str,
        state: Annotated[dict, InjectedState] = None,
) -> str:
    """Export a validated golden-dataset JSON payload to CSV and run quality metrics.

    For the ``golden-dataset`` target, call this after ``render_target_output`` with the
    same JSON. This writes the CSV under the output folder and runs metrics in one step.
    It also generates:
    - `/golden_dataset_metrics.md`: Markdown table of all items with quality metrics
    - `/final_report.md`: Comprehensive report of the entire golden dataset generation process

    Args:
        payload_json: JSON object matching the golden-dataset schema (same payload as ``render_target_output``).
        state: LangGraph state

    Returns:
        Paths to the exported CSV and metrics output, or a validation or runtime error message.
    """
    import time
    from research_agent.skills.golden_dataset.pipeline import (
        GOLDEN_DATASET_TARGET_ID,
        evaluate_and_report_golden_dataset,
        export_golden_dataset_csv,
    )

    _, payload, err = _prepare_validated_payload(GOLDEN_DATASET_TARGET_ID, payload_json)
    if err:
        return err

    try:
        # Calculate elapsed time from agent start
        start_time = state.get("agent_start_time") if state else None
        elapsed_seconds = (time.time() - start_time) if start_time else 0.0

        output_subfolder = Path(os.environ.get("OUTPUT_FOLDER", REPORTS_OUTPUT_FOLDER))
        csv_path = export_golden_dataset_csv(payload, output_subfolder)

        # Evaluate and generate reports
        metrics_csv_path, markdown_content, final_report_content = evaluate_and_report_golden_dataset(
            csv_path=csv_path,
            payload=payload,
            output_folder=output_subfolder,
            elapsed_seconds=elapsed_seconds
        )

        # Write metrics markdown file
        metrics_md_path = output_subfolder / "golden_dataset_metrics.md"
        with open(metrics_md_path, "w", encoding="utf-8") as f:
            f.write(markdown_content)

        # Write final report
        final_report_path = output_subfolder / "final_report.md"
        with open(final_report_path, "w", encoding="utf-8") as f:
            f.write(final_report_content)

        # Update state files if available
        if state is not None:
            files = state.get("files", {})
            files["/golden_dataset_metrics.md"] = create_file_data(markdown_content)
            files["/final_report.md"] = create_file_data(final_report_content)
            state["files"] = files

        return (
            f"**CSV exported to:** `{csv_path}`\n\n"
            f"**Metrics CSV:** `{metrics_csv_path}`\n\n"
            f"**Metrics Markdown:** `{metrics_md_path}`\n\n"
            f"**Final Report:** `{final_report_path}`\n\n"
            f"All files have been generated successfully!"
        )
    except Exception as e:
        return f"**Error exporting or evaluating golden dataset:** {e}"


@tool(parse_docstring=True)
def trigger_dataset_evaluation(file_path: str) -> str:
    """Evaluate a generated golden dataset CSV to compute quality metrics.

    Run this tool only after you have successfully generated a golden dataset
    and received the CSV file path locally in your output folder. This runs a heavy
    evaluation script to compute Similarity, Relevance, Coherence, and Groundedness.

    For new datasets, prefer ``finalize_golden_dataset_output``, which exports the CSV
    and runs this evaluation in order. Use this tool to re-run metrics on an existing CSV.

    Args:
        file_path: The path to the CSV file to evaluate (e.g., "./output/golden_dataset.csv").

    Returns:
        The result of the quality metric evaluation, including the path to the scored dataset.
    """
    from research_agent.skills.golden_dataset.pipeline import evaluate_golden_dataset_csv_file

    return evaluate_golden_dataset_csv_file(file_path)
