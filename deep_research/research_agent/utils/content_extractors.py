from __future__ import annotations

from pathlib import Path

import pymupdf4llm
import pypdf
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from logger_utils import setup_logger

logger = setup_logger(__name__)


def _extract_pdf_text(file_path: Path) -> str:
    """Extract PDF content as markdown without ML model downloads.\n\n    Returns:\n        str: Extracted content from the PDF file.\n    """
    try:
        logger.info("Use PyMuPDF4LLM for PDF markdown extraction.")

        markdown_content = pymupdf4llm.to_markdown(str(file_path))
        if isinstance(markdown_content, list):
            # Convert each dictionary to a string representation
            return "\n\n".join(str(item) for item in markdown_content)
        if markdown_content.strip():
            return markdown_content
    except Exception as e:
        logger.error(f"PyMuPDF4LLM PDF extraction failed: {e}")
        # Fallback to pypdf if markdown extraction fails
        try:
            logger.info("Falling back to pypdf for PDF text extraction.")
            reader = pypdf.PdfReader(file_path)
            page_texts: list[str] = []
            for index, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    page_texts.append(f"## Page {index}:\n\n{text}")
            return "\n\n".join(page_texts)
        except Exception as e:
            return f"Error extracting PDF text: {e}"
    return ""


def _extract_text_file(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8")


def _extract_docx_text(file_path: Path) -> str:
    document = Document(str(file_path))
    paragraphs = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text and paragraph.text.strip()
    ]

    # Many documents are table-based; include table cells so content is not silently missed.
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                call_paragraph = [
                    paragraph.text.strip()
                    for paragraph in cell.paragraphs
                    if paragraph.text and paragraph.text.strip()
                ]
                if call_paragraph:
                    paragraphs.extend(call_paragraph)

    return "\n".join(paragraphs)


def _extract_pptx_text(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    slide_sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide:
            notes = []
            for shape in slide.notes_slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    notes.append(text.strip())
            notes_text = "\n".join(notes)
        if notes_text:
            parts.append(f"Speaker Notes:\n{notes_text}")
        slide_sections.append("\n".join(parts))

    return "\n\n".join(slide_sections)


def _extract_xlsx_text(file_path: Path) -> str:
    workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    sections: list[str] = []
    try:
        for worksheet in workbook.worksheets:
            rows: list[str] = []
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if values:
                    rows.append(" | ".join(values))
            body = "\n".join(rows) if rows else "(empty sheet)"
            sections.append(f"Sheet: {worksheet.title}\n{body}")
    finally:
        workbook.close()

    return "\n\n".join(sections)


def extract_supported_document(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf_text(file_path)
    if suffix in {".txt", ".md"}:
        return _extract_text_file(file_path)
    if suffix == ".docx":
        return _extract_docx_text(file_path)
    if suffix == ".pptx":
        return _extract_pptx_text(file_path)
    if suffix == ".xlsx":
        return _extract_xlsx_text(file_path)
    raise ValueError(f"Unsupported document type: {suffix}")
